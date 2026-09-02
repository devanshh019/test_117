# PowerPoint Presentation (.pptx) Widescreen Strategy Renderer
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from .base import Renderer
from ..schemas import PptxSpec, DocumentResult, DocType


class PptxRenderer(Renderer):

    def render(self, spec: PptxSpec, output_path: Path) -> DocumentResult:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        self._title_slide(prs, spec.title, spec.subtitle or "")
        for slide in spec.slides:
            self._content_slide(prs, slide.title, slide.bullets)

        prs.save(output_path)

        return DocumentResult(
            filename=output_path.name,
            path=str(output_path),
            doc_type=DocType.PPTX,
            size_bytes=os.path.getsize(output_path),
        )

    def _title_slide(self, prs: Presentation, title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*self.theme.rgb_dark)

        box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.9), Inches(3.5))
        tf = box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title.upper()
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(self.theme.title_size)
        tf.paragraphs[0].font.color.rgb = RGBColor(*self.theme.rgb_soft_bg)

        if subtitle:
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*self.theme.rgb_orange)

    def _content_slide(self, prs: Presentation, title: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*self.theme.rgb_soft_bg)

        hdr = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
        hdr.text_frame.paragraphs[0].text = title
        hdr.text_frame.paragraphs[0].font.bold = True
        hdr.text_frame.paragraphs[0].font.size = Pt(self.theme.header_size)
        hdr.text_frame.paragraphs[0].font.color.rgb = RGBColor(*self.theme.rgb_dark)

        body = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.8))
        tf = body.text_frame
        tf.word_wrap = True
        for b in bullets:
            p = tf.add_paragraph()
            p.text = f"•  {b}"
            p.font.size = Pt(self.theme.body_size)
            p.font.color.rgb = RGBColor(*self.theme.rgb_muted)
