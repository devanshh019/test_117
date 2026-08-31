#UDIT

import os
import re
import time
from pathlib import Path
from typing import Dict, Any, List, Optional

import docx
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor

from .config import STORAGE_DIR, THEME


class IndustrialDocumentGenerator:
    """Generates styled Office documents (.pptx, .docx, .xlsx) from structured outputs."""

    def __init__(self):
        self.storage_dir = Path(STORAGE_DIR)
        self.storage_dir.mkdir(parents=True, exist_ok=True)

    def _get_output_path(self, prefix: str, raw_title: str, ext: str) -> tuple[str, Path]:
        """Generates a clean timestamped filename and absolute path."""
        clean_title = re.sub(r'[^a-zA-Z0-9]', '_', raw_title)[:24]
        filename = f"{prefix}_{clean_title}_{int(time.time())}.{ext}"
        filepath = self.storage_dir / filename
        return filename, filepath


    # PowerPoint Generation Helpers
 
    def _add_pptx_title_slide(self, prs: Presentation, title: str, subtitle: str):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        bg = slide.shapes.add_shape(1, 0, 0, PptxInches(13.333), PptxInches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = PptxRGBColor(*THEME["rgb_dark"])

        box = slide.shapes.add_textbox(PptxInches(1.2), PptxInches(2.0), PptxInches(10.9), PptxInches(3.5))
        tf = box.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title.upper()
        p0.font.bold = True
        p0.font.size = PptxPt(THEME["title_size"])
        p0.font.color.rgb = PptxRGBColor(*THEME["rgb_soft_bg"])

        p1 = tf.add_paragraph()
        p1.text = subtitle or THEME["subtitle"]
        p1.font.size = PptxPt(18)
        p1.font.color.rgb = PptxRGBColor(*THEME["rgb_orange"])
        p1.space_before = PptxPt(12)

        p2 = tf.add_paragraph()
        p2.text = THEME["confidential_tag"]
        p2.font.size = PptxPt(11)
        p2.font.color.rgb = PptxRGBColor(*THEME["rgb_subtle"])
        p2.space_before = PptxPt(10)

    def _add_pptx_content_slide(self, prs: Presentation, slide_info: Dict[str, Any], fallback_title: str):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        bg = slide.shapes.add_shape(1, 0, 0, PptxInches(13.333), PptxInches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = PptxRGBColor(*THEME["rgb_soft_bg"])

        # Header
        hdr_box = slide.shapes.add_textbox(PptxInches(1.0), PptxInches(0.8), PptxInches(11.3), PptxInches(1.0))
        tf_hdr = hdr_box.text_frame
        p_hdr = tf_hdr.paragraphs[0]
        p_hdr.text = slide_info.get("title") or fallback_title
        p_hdr.font.bold = True
        p_hdr.font.size = PptxPt(THEME["header_size"])
        p_hdr.font.color.rgb = PptxRGBColor(*THEME["rgb_dark"])

        # Bullets
        body_box = slide.shapes.add_textbox(PptxInches(1.0), PptxInches(2.0), PptxInches(11.3), PptxInches(4.8))
        tf_body = body_box.text_frame
        tf_body.word_wrap = True

        for b in slide_info.get("bullets", []):
            p = tf_body.add_paragraph()
            p.text = f"•  {b}"
            p.font.size = PptxPt(THEME["body_size"])
            p.font.color.rgb = PptxRGBColor(*THEME["rgb_muted"])
            p.space_after = PptxPt(14)

    def generate_custom_powerpoint(
        self,
        presentation_title: str,
        subtitle: str = None,
        slides_data: List[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        """Creates a 16:9 widescreen PowerPoint deck."""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        subtitle_text = subtitle or THEME["subtitle"]
        self._add_pptx_title_slide(prs, presentation_title, subtitle_text)

        slides = slides_data or []
        for s_info in slides:
            self._add_pptx_content_slide(prs, s_info, presentation_title)

        filename, filepath = self._get_output_path("Presentation", presentation_title, "pptx")
        prs.save(filepath)

        return {
            "filename": filename,
            "path": f"/api/artifacts/{filename}",
            "file_type": "pptx",
            "size_bytes": os.path.getsize(filepath),
            "title": f"{presentation_title} (.pptx)",
            "slides": [{"title": presentation_title, "bullets": [subtitle_text]}] + slides,
        }


    # Word Document Generation Helpers

    def generate_custom_word_doc(
        self,
        doc_title: str,
        subject: str,
        paragraphs: List[str] = None,
    ) -> Dict[str, Any]:
        """Creates a Word document report with styled headers and paragraphs."""
        doc = docx.Document()
        style = doc.styles["Normal"]
        font = style.font
        font.name = THEME["font_family"]
        font.size = Pt(THEME["doc_body_size"])
        font.color.rgb = RGBColor(*THEME["rgb_dark"])

        # Organization Header
        p_hdr = doc.add_paragraph()
        p_hdr.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run_org = p_hdr.add_run(f"{THEME['org_title']}\n")
        run_org.bold = True
        run_org.font.size = Pt(THEME["doc_title_size"])
        run_org.font.color.rgb = RGBColor(*THEME["rgb_orange"])

        run_title = p_hdr.add_run(f"{doc_title.upper()}\n")
        run_title.bold = True
        run_title.font.size = Pt(THEME["doc_heading_size"])

        # Subject line
        p_subj = doc.add_paragraph()
        r_lbl = p_subj.add_run("SUBJECT: ")
        r_lbl.bold = True
        r_lbl.font.color.rgb = RGBColor(*THEME["rgb_orange"])
        p_subj.add_run(subject)

        # Body paragraphs
        content = paragraphs or []
        for idx, text in enumerate(content, 1):
            h = doc.add_heading(f"Section {idx}: Technical Assessment", level=2)
            h.runs[0].font.size = Pt(THEME["doc_heading_size"])
            h.runs[0].font.color.rgb = RGBColor(*THEME["rgb_dark"])
            doc.add_paragraph(text)

        filename, filepath = self._get_output_path("Document", doc_title, "docx")
        doc.save(filepath)

        return {
            "filename": filename,
            "path": f"/api/artifacts/{filename}",
            "file_type": "docx",
            "size_bytes": os.path.getsize(filepath),
            "title": f"{doc_title} (.docx)",
        }

    # -------------------------------------------------------------------------
    # Excel Spreadsheet Generation Helpers
    # -------------------------------------------------------------------------
    def generate_custom_excel(
        self,
        sheet_title: str = "Engineering_Data",
        headers: List[str] = None,
        rows_data: List[List[Any]] = None,
    ) -> Dict[str, Any]:
        """Creates an Excel spreadsheet with styled header row and bordered data."""
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Assessment_Data"
        ws.views.sheetView[0].showGridLines = True

        # Title Banner (Row 1)
        ws.merge_cells("A1:H1")
        ws["A1"] = sheet_title.upper()
        ws["A1"].font = Font(name=THEME["font_family"], size=13, bold=True, color=THEME["text_white"])
        ws["A1"].fill = PatternFill(start_color=THEME["primary_orange"], end_color=THEME["primary_orange"], fill_type="solid")
        ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
        ws.row_dimensions[1].height = 28

        # Table Column Headers (Row 3)
        cols = headers or ["Item ID", "Parameter", "Calculated Value", "Unit", "Status"]
        ws.row_dimensions[3].height = 24

        for col_num, h in enumerate(cols, 1):
            cell = ws.cell(row=3, column=col_num, value=h)
            cell.font = Font(name=THEME["font_family"], size=10, bold=True, color=THEME["text_white"])
            cell.fill = PatternFill(start_color=THEME["dark_navy"], end_color=THEME["dark_navy"], fill_type="solid")
            cell.alignment = Alignment(horizontal="center", vertical="center")

        # Table Data Rows
        thin_border = Border(
            left=Side(style="thin", color=THEME["border_light"]),
            right=Side(style="thin", color=THEME["border_light"]),
            top=Side(style="thin", color=THEME["border_light"]),
            bottom=Side(style="thin", color=THEME["border_light"]),
        )

        data = rows_data or []
        for r_idx, row_values in enumerate(data, 4):
            ws.row_dimensions[r_idx].height = 20
            for c_idx, val in enumerate(row_values, 1):
                cell = ws.cell(row=r_idx, column=c_idx, value=val)
                cell.font = Font(name=THEME["font_family"], size=9.5)
                cell.border = thin_border
                cell.alignment = Alignment(horizontal="center", vertical="center")

        # Auto-fit column widths
        for col in ws.columns:
            max_len = max(len(str(cell.value or "")) for cell in col)
            col_letter = openpyxl.utils.get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = max(max_len + 4, 14)

        filename, filepath = self._get_output_path("Spreadsheet", sheet_title, "xlsx")
        wb.save(filepath)

        return {
            "filename": filename,
            "path": f"/api/artifacts/{filename}",
            "file_type": "xlsx",
            "size_bytes": os.path.getsize(filepath),
            "title": f"{sheet_title} (.xlsx)",
        }


# Shared document generator instance
doc_generator = IndustrialDocumentGenerator()

